#!/usr/bin/env python3
"""Generate a PPT summary for a journal volume or issue page URL.

Features:
- Search Crossref for a journal by title + volume
- Or fetch a supplied issue/article-list URL directly and extract titles/authors from the page
- Keep the first N titles
- Optionally translate titles to Chinese (best effort via Google public endpoint)
- Generate a PPTX using the same visual style as the existing Automatica deck
- Save fetched metadata to JSON for reuse/auditing

Examples:
  python3 make_journal_ppt.py --journal Automatica --volume 185 --count 50
  python3 make_journal_ppt.py --url "https://www.sciencedirect.com/journal/automatica/vol/185" --count 50
  python3 make_journal_ppt.py --journal Automatica --volume 185 --count 10 --no-translate
"""

from __future__ import annotations

import argparse
import base64
import hashlib
import html
import json
import math
import os
import random
import re
import socket
import subprocess
import sys
import tempfile
import time
import urllib.parse
import urllib.request
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

CROSSREF_API = "https://api.crossref.org"
USER_AGENT = "journal-ppt-generator/1.0 (OpenClaw workspace script)"

NAVY = RGBColor(28, 43, 74)
TEAL = RGBColor(46, 95, 117)
LIGHT = RGBColor(246, 248, 251)
MID = RGBColor(103, 118, 136)
DARK = RGBColor(35, 42, 52)
ACCENT = RGBColor(214, 228, 240)
WHITE = RGBColor(255, 255, 255)

TITLE_FONT = 14.0
AUTHOR_FONT = 10.5
ZH_FONT = 11.0
CONTENT_TOP = 1.0
CONTENT_BOTTOM = 7.18
CONTENT_HEIGHT = CONTENT_BOTTOM - CONTENT_TOP
CARD_X = 0.65
CARD_W = 12.0
NUM_X = 0.82
LEFT = 1.48
TEXT_W = 10.75
CARD_GAP = 0.14

LAYOUT5 = {
    "card_h": 1.08,
    "title_y": 0.08,
    "title_h": 0.27,
    "author_y": 0.6,
    "author_h": 0.18,
    "zh_y": 0.77,
    "zh_h": 0.18,
    "title_chars": 82,
    "author_chars": 105,
    "zh_chars": 84,
}

LAYOUT4 = {
    "card_h": 1.34,
    "title_y": 0.09,
    "title_h": 0.42,
    "author_y": 0.62,
    "author_h": 0.20,
    "zh_y": 0.83,
    "zh_h": 0.26,
    "title_chars": 82,
    "author_chars": 105,
    "zh_chars": 84,
}


class FetchError(RuntimeError):
    pass


class SimpleWebSocket:
    def __init__(self, ws_url: str, timeout: int = 20):
        self.ws_url = ws_url
        self.timeout = timeout
        self.sock: socket.socket | None = None

    def connect(self) -> None:
        parsed = urllib.parse.urlparse(self.ws_url)
        host = parsed.hostname or "127.0.0.1"
        port = parsed.port or 80
        path = parsed.path or "/"
        if parsed.query:
            path += "?" + parsed.query
        sock = socket.create_connection((host, port), timeout=self.timeout)
        key = base64.b64encode(os.urandom(16)).decode("ascii")
        req = (
            f"GET {path} HTTP/1.1\r\n"
            f"Host: {host}:{port}\r\n"
            "Upgrade: websocket\r\n"
            "Connection: Upgrade\r\n"
            f"Sec-WebSocket-Key: {key}\r\n"
            "Sec-WebSocket-Version: 13\r\n\r\n"
        )
        sock.sendall(req.encode("utf-8"))
        resp = self._recv_until(sock, b"\r\n\r\n")
        if b" 101 " not in resp.split(b"\r\n", 1)[0]:
            raise FetchError(f"WebSocket 握手失败：{resp[:200]!r}")
        self.sock = sock

    def close(self) -> None:
        if self.sock is not None:
            try:
                self.sock.close()
            finally:
                self.sock = None

    def send_json(self, payload: dict[str, Any]) -> None:
        raw = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        self._send_frame(raw)

    def recv_json(self) -> dict[str, Any]:
        data = self._recv_frame()
        return json.loads(data.decode("utf-8"))

    def _recv_until(self, sock: socket.socket, marker: bytes) -> bytes:
        chunks = []
        data = b""
        while marker not in data:
            chunk = sock.recv(4096)
            if not chunk:
                break
            chunks.append(chunk)
            data = b"".join(chunks)
        return data

    def _send_frame(self, payload: bytes) -> None:
        if self.sock is None:
            raise FetchError("WebSocket 未连接")
        first = 0x81
        mask_bit = 0x80
        length = len(payload)
        header = bytearray([first])
        if length < 126:
            header.append(mask_bit | length)
        elif length < (1 << 16):
            header.append(mask_bit | 126)
            header.extend(length.to_bytes(2, "big"))
        else:
            header.append(mask_bit | 127)
            header.extend(length.to_bytes(8, "big"))
        mask = os.urandom(4)
        header.extend(mask)
        masked = bytes(b ^ mask[i % 4] for i, b in enumerate(payload))
        self.sock.sendall(bytes(header) + masked)

    def _recv_exact(self, n: int) -> bytes:
        if self.sock is None:
            raise FetchError("WebSocket 未连接")
        chunks = []
        remaining = n
        while remaining > 0:
            chunk = self.sock.recv(remaining)
            if not chunk:
                raise FetchError("WebSocket 连接意外关闭")
            chunks.append(chunk)
            remaining -= len(chunk)
        return b"".join(chunks)

    def _recv_frame(self) -> bytes:
        header = self._recv_exact(2)
        b1, b2 = header[0], header[1]
        opcode = b1 & 0x0F
        masked = (b2 & 0x80) != 0
        length = b2 & 0x7F
        if length == 126:
            length = int.from_bytes(self._recv_exact(2), "big")
        elif length == 127:
            length = int.from_bytes(self._recv_exact(8), "big")
        mask = self._recv_exact(4) if masked else b""
        payload = self._recv_exact(length) if length else b""
        if masked:
            payload = bytes(b ^ mask[i % 4] for i, b in enumerate(payload))
        if opcode == 0x8:
            self.close()
            return b""
        if opcode == 0x9:
            # ping -> pong
            pong = bytearray([0x8A, len(payload)]) + payload
            assert self.sock is not None
            self.sock.sendall(bytes(pong))
            return self._recv_frame()
        return payload


def http_get_json(url: str, timeout: int = 30) -> dict[str, Any]:
    req = urllib.request.Request(url, headers={"User-Agent": USER_AGENT})
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        return json.loads(resp.read().decode("utf-8"))


def http_get_text(url: str, timeout: int = 30) -> str:
    req = urllib.request.Request(url, headers={"User-Agent": USER_AGENT})
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        return resp.read().decode("utf-8")


def slugify(text: str) -> str:
    text = text.strip().lower()
    text = re.sub(r"[^a-z0-9]+", "_", text)
    text = re.sub(r"_+", "_", text).strip("_")
    return text or "journal"


def normalize_title(title: str) -> str:
    return re.sub(r"\s+", " ", title).strip()


def best_journal_match(journal_title: str) -> dict[str, Any]:
    q = urllib.parse.quote(journal_title)
    url = f"{CROSSREF_API}/journals?query={q}&rows=10"
    data = http_get_json(url)
    items = data.get("message", {}).get("items", [])
    if not items:
        raise FetchError(f"未找到期刊：{journal_title}")

    target = journal_title.strip().lower()
    exact = []
    partial = []
    for item in items:
        title = (item.get("title") or "").strip()
        title_lower = title.lower()
        if title_lower == target:
            exact.append(item)
        elif target in title_lower or title_lower in target:
            partial.append(item)
    return (exact or partial or items)[0]


def parse_authors(author_list: list[dict[str, Any]]) -> str:
    names = []
    for a in author_list or []:
        given = (a.get("given") or "").strip()
        family = (a.get("family") or "").strip()
        name = " ".join(x for x in [given, family] if x)
        if not name:
            name = (a.get("name") or "").strip()
        if name:
            names.append(name)
    return ", ".join(names) if names else "Unknown"


def parse_author_names(author_value: Any) -> str:
    if isinstance(author_value, list):
        names = []
        for a in author_value:
            if isinstance(a, dict):
                name = (a.get("name") or "").strip()
                if not name:
                    given = (a.get("given") or "").strip()
                    family = (a.get("family") or "").strip()
                    name = " ".join(x for x in [given, family] if x).strip()
                if name:
                    names.append(name)
            elif isinstance(a, str) and a.strip():
                names.append(a.strip())
        return ", ".join(names) if names else "Unknown"
    if isinstance(author_value, dict):
        return parse_author_names([author_value])
    if isinstance(author_value, str) and author_value.strip():
        return author_value.strip()
    return "Unknown"


def extract_meta_content(html_text: str, key: str) -> str:
    patterns = [
        rf'<meta[^>]+name=["\']{re.escape(key)}["\'][^>]+content=["\']([^"\']+)["\']',
        rf'<meta[^>]+property=["\']{re.escape(key)}["\'][^>]+content=["\']([^"\']+)["\']',
        rf'<meta[^>]+content=["\']([^"\']+)["\'][^>]+name=["\']{re.escape(key)}["\']',
        rf'<meta[^>]+content=["\']([^"\']+)["\'][^>]+property=["\']{re.escape(key)}["\']',
    ]
    for pattern in patterns:
        m = re.search(pattern, html_text, flags=re.I)
        if m:
            return html.unescape(m.group(1)).strip()
    return ""


def extract_title_tag(html_text: str) -> str:
    m = re.search(r"<title>(.*?)</title>", html_text, flags=re.I | re.S)
    if not m:
        return ""
    return html.unescape(re.sub(r"\s+", " ", m.group(1))).strip()


def guess_journal_and_volume_from_html(page_url: str, html_text: str) -> tuple[str, str]:
    candidates = [
        extract_meta_content(html_text, "citation_journal_title"),
        extract_meta_content(html_text, "prism.publicationName"),
        extract_meta_content(html_text, "og:site_name"),
        extract_title_tag(html_text),
        page_url,
    ]
    journal = ""
    volume = ""
    for text in candidates:
        if not text:
            continue
        vm = re.search(r"\bVolume\s+(\d+)\b", text, flags=re.I)
        if vm and not volume:
            volume = vm.group(1)
        if not journal:
            for part in re.split(r"[|·—-]", text):
                part = part.strip()
                if part and "volume" not in part.lower() and len(part) < 120:
                    journal = part
                    break
    return journal or "Journal", volume or "Unknown"


def _walk_json_for_article_like_entries(node: Any, out: list[dict[str, str]]) -> None:
    if isinstance(node, dict):
        title = node.get("headline") or node.get("name") or node.get("title")
        authors = node.get("author") or node.get("authors") or node.get("creator")
        if isinstance(title, str) and title.strip():
            title_clean = normalize_title(html.unescape(title))
            if title_clean and len(title_clean) > 8:
                out.append({
                    "title": title_clean,
                    "authors": parse_author_names(authors),
                    "title_zh": "",
                    "source": "URL",
                })
        for v in node.values():
            _walk_json_for_article_like_entries(v, out)
    elif isinstance(node, list):
        for item in node:
            _walk_json_for_article_like_entries(item, out)


SCIENCEDIRECT_TITLE_BLOCK = re.compile(
    r'<span[^>]*class=["\'][^"\']*title-text[^"\']*["\'][^>]*>(.*?)</span>',
    flags=re.I | re.S,
)
SCIENCEDIRECT_AUTHOR_BLOCK = re.compile(
    r'<span[^>]*class=["\'][^"\']*author[^"\']*["\'][^>]*>(.*?)</span>',
    flags=re.I | re.S,
)
IEEE_HEADING_LINK_BLOCK = re.compile(
    r'<h[23][^>]*>\s*<a[^>]+href=["\']([^"\']*(?:/document/\d+|[?&]arnumber=\d+)[^"\']*)["\'][^>]*>(.*?)</a>\s*</h[23]>',
    flags=re.I | re.S,
)
IEEE_RESULT_LINK_BLOCK = re.compile(
    r'<a[^>]+class=["\'][^"\']*(?:art-abs-url|result-item-title-link)[^"\']*["\'][^>]+href=["\']([^"\']*(?:/document/\d+|[?&]arnumber=\d+)[^"\']*)["\'][^>]*>(.*?)</a>',
    flags=re.I | re.S,
)


def strip_tags(text: str) -> str:
    text = re.sub(r"<[^>]+>", " ", text)
    text = html.unescape(text)
    return normalize_title(text)


def titleize_slug(slug: str) -> str:
    parts = [p for p in re.split(r"[-_]+", slug) if p]
    if not parts:
        return slug
    keep_lower = {"of", "and", "the", "in", "on", "for", "a", "an"}
    titled = []
    for i, p in enumerate(parts):
        if i > 0 and p.lower() in keep_lower:
            titled.append(p.lower())
        else:
            titled.append(p.capitalize())
    return " ".join(titled)


def infer_journal_volume_from_url(page_url: str) -> tuple[str, str]:
    parsed = urllib.parse.urlparse(page_url)
    path = parsed.path

    # ScienceDirect issue pages: /journal/automatica/vol/185/suppl/C
    m = re.search(r"/journal/([^/]+)/vol/(\d+)", path, flags=re.I)
    if m:
        return titleize_slug(m.group(1)), m.group(2)

    # Generic volume hints in URL
    vm = re.search(r"(?:^|[/_-])vol(?:ume)?[/_-]?(\d+)(?:$|[/_-])", path, flags=re.I)
    volume = vm.group(1) if vm else ""

    parts = [p for p in path.split("/") if p]
    journal = ""
    for part in parts:
        low = part.lower()
        if low in {"journal", "journals", "issue", "issues", "vol", "volume", "suppl", "supplement"}:
            continue
        if part.isdigit():
            continue
        if len(part) > 2 and re.search(r"[a-zA-Z]", part):
            journal = titleize_slug(part)
            break
    return journal or "Journal", volume or "Unknown"


def extract_ieee_issue_entries_from_html(html_text: str, count: int) -> list[dict[str, str]]:
    candidates: list[dict[str, str]] = []
    for _, raw_title in IEEE_HEADING_LINK_BLOCK.findall(html_text):
        title = strip_tags(raw_title)
        if title and len(title) > 8:
            candidates.append({"title": title, "authors": "Unknown", "title_zh": "", "source": "IEEE-HTML"})
            if len(candidates) >= count:
                return candidates
    if len(candidates) >= count:
        return candidates

    seen = {item["title"] for item in candidates}
    for _, raw_title in IEEE_RESULT_LINK_BLOCK.findall(html_text):
        title = strip_tags(raw_title)
        if not title or len(title) <= 8 or title in seen:
            continue
        seen.add(title)
        candidates.append({"title": title, "authors": "Unknown", "title_zh": "", "source": "IEEE-HTML"})
        if len(candidates) >= count:
            break
    return candidates


def parse_ieee_issue_params(page_url: str) -> dict[str, str]:
    parsed = urllib.parse.urlparse(page_url)
    qs = urllib.parse.parse_qs(parsed.query)
    return {
        "publication_number": (qs.get("punumber") or [""])[0],
        "issue_number": (qs.get("isnumber") or [""])[0],
        "sort_type": (qs.get("sortType") or [""])[0],
        "rows_per_page": (qs.get("rowsPerPage") or [""])[0],
        "page_number": (qs.get("pageNumber") or ["1"])[0],
    }


def extract_ieee_items_via_api(page_url: str, count: int, wait_s: float = 8.0) -> tuple[dict[str, Any], list[dict[str, str]]]:
    params = parse_ieee_issue_params(page_url)
    pub = params["publication_number"]
    issue = params["issue_number"]
    if not pub or not issue:
        raise FetchError(f"IEEE issue URL is missing punumber/isnumber: {page_url}")

    rows = max(count, int(params["rows_per_page"] or "25"))
    start_page = max(1, int(params["page_number"] or "1"))
    sort_type = params["sort_type"]

    ws_url, launched_proc, launched_profile = cdp_get_or_create_page_ws_auto(
        page_url,
        auto_launch=True,
        headless=False,
    )
    ws = SimpleWebSocket(ws_url)
    ws.connect()
    try:
        cid = 1
        cdp_call(ws, cid, "Page.enable"); cid += 1
        cdp_call(ws, cid, "Runtime.enable"); cid += 1
        cdp_call(ws, cid, "Page.navigate", {"url": page_url}, wait_s=20.0); cid += 1
        time.sleep(10.0)
        consent_expr = """
(() => {
  const norm = (s) => (s || '').replace(/\\s+/g, ' ').trim().toLowerCase();
  const buttons = Array.from(document.querySelectorAll('button, a, [role="button"], input[type="button"]'));
  const match = buttons.find((el) => {
    const text = norm(el.innerText || el.textContent || el.value || '');
    return (
      text === 'accept' ||
      text === 'i accept' ||
      text.includes('accept all') ||
      text.includes('accept cookies') ||
      text === '接受' ||
      text.includes('接受全部')
    );
  });
  if (match) {
    match.click();
    return true;
  }
  return false;
})()
"""
        cdp_call(ws, cid, "Runtime.evaluate", {"expression": consent_expr, "returnByValue": True}, wait_s=10.0); cid += 1
        time.sleep(4.0)

        fetch_expr = f"""
(async () => {{
  const banned = /table of contents|publication information|society information/i;
  const pub = {json.dumps(pub)};
  const issue = {json.dumps(issue)};
  const rows = {rows};
  const startPage = {start_page};
  const needed = {count};
  const sortType = {json.dumps(sort_type)};
  const meta = await fetch(`/rest/publication/${{pub}}/issue/${{issue}}/metadata`, {{
    credentials: 'include'
  }}).then(r => r.json());

  const items = [];
  let totalRecords = 0;
  let totalPages = 1;
  let page = startPage;
  while (items.length < needed && page <= totalPages) {{
    const body = {{ pageNumber: page, rowsPerPage: rows }};
    if (sortType) body.sortType = sortType;
    const toc = await fetch(`/rest/search/pub/${{pub}}/issue/${{issue}}/toc`, {{
      method: 'POST',
      credentials: 'include',
      headers: {{ 'Content-Type': 'application/json' }},
      body: JSON.stringify(body),
    }}).then(r => r.json());

    totalRecords = toc.totalRecords || totalRecords || 0;
    totalPages = toc.totalPages || totalPages || 1;
    const records = Array.isArray(toc.records) ? toc.records : [];
    for (const rec of records) {{
      const title = (rec.articleTitle || '').replace(/\\s+/g, ' ').trim();
      if (!title || banned.test(title)) continue;
      const authors = Array.isArray(rec.authors)
        ? rec.authors.map(a => (a.preferredName || a.name || '')).filter(Boolean).join('; ')
        : '';
      items.push({{
        title,
        authors: authors || 'Unknown',
        title_zh: '',
        source: 'IEEE-API',
      }});
      if (items.length >= needed) break;
    }}
    page += 1;
  }}

  return {{
    meta: {{
      title: meta.publicationTitle || 'IEEE Journal',
      volume: meta.volume || 'Unknown',
      issue: meta.issue || '',
      year: meta.year || '',
      month: meta.month || '',
      totalRecords,
      totalPages,
    }},
    items,
    pageTitle: document.title || '',
  }};
}})()
"""
        result = cdp_call(
            ws,
            cid,
            "Runtime.evaluate",
            {"expression": fetch_expr, "returnByValue": True, "awaitPromise": True},
            wait_s=max(40.0, wait_s * 5),
        )
        value = (((result.get("result") or {}).get("result") or {}).get("value"))
        if not value or not value.get("items"):
            raise FetchError(f"IEEE API extraction failed for URL: {page_url} | page_title={(value or {}).get('pageTitle', '')!r}")
        meta = value.get("meta") or {}
        items = [
            {
                "title": normalize_title(it.get("title", "")),
                "authors": normalize_title(it.get("authors", "")) or "Unknown",
                "title_zh": "",
                "source": "IEEE-API",
            }
            for it in value.get("items", [])
            if normalize_title(it.get("title", ""))
        ]
        return {
            "title": meta.get("title") or "IEEE Journal",
            "volume": str(meta.get("volume") or "Unknown"),
            "issue": str(meta.get("issue") or ""),
            "year": str(meta.get("year") or ""),
            "month": str(meta.get("month") or ""),
            "source_url": page_url,
            "total_records": meta.get("totalRecords") or len(items),
            "total_pages": meta.get("totalPages") or 1,
        }, items[:count]
    finally:
        ws.close()
        if launched_proc is not None:
            try:
                launched_proc.terminate()
                launched_proc.wait(timeout=5)
            except Exception:
                pass
        if launched_profile:
            try:
                import shutil
                shutil.rmtree(launched_profile, ignore_errors=True)
            except Exception:
                pass


def find_browser_executable() -> str:
    candidates = [
        r"C:\Program Files\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
        r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
        r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
    ]
    for path in candidates:
        if os.path.exists(path):
            return path
    raise FetchError("No Chrome or Edge executable found for browser extraction")


def start_debug_browser(port: int = 18800, timeout: int = 20, headless: bool = True) -> tuple[subprocess.Popen[bytes], str]:
    browser = find_browser_executable()
    profile_dir = tempfile.mkdtemp(prefix="journal_ppt_cdp_")
    cmd = [
        browser,
        f"--remote-debugging-port={port}",
        "--disable-gpu",
        "--no-first-run",
        "--no-default-browser-check",
        "--disable-blink-features=AutomationControlled",
        f"--user-data-dir={profile_dir}",
        "about:blank",
    ]
    if headless:
        cmd.insert(2, "--headless=new")
    proc = subprocess.Popen(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    deadline = time.time() + timeout
    last_err: Exception | None = None
    while time.time() < deadline:
        try:
            http_get_json(f"http://127.0.0.1:{port}/json/version", timeout=2)
            return proc, profile_dir
        except Exception as exc:
            last_err = exc
            time.sleep(0.5)
    try:
        proc.terminate()
    except Exception:
        pass
    raise FetchError(f"Failed to start browser debugging endpoint: {last_err}")


def cdp_get_or_create_page_ws_auto(
    page_url: str,
    timeout: int = 20,
    auto_launch: bool = False,
    headless: bool = True,
) -> tuple[str, subprocess.Popen[bytes] | None, str | None]:
    launched_proc: subprocess.Popen[bytes] | None = None
    launched_profile: str | None = None
    try:
        tabs = http_get_json("http://127.0.0.1:18800/json/list", timeout=timeout)
    except Exception:
        if not auto_launch:
            raise
        launched_proc, launched_profile = start_debug_browser(port=18800, timeout=timeout, headless=headless)
        tabs = http_get_json("http://127.0.0.1:18800/json/list", timeout=timeout)
    if isinstance(tabs, list):
        for tab in tabs:
            tab_url = tab.get("url") or ""
            if tab_url == page_url or tab_url.startswith(page_url):
                ws_url = tab.get("webSocketDebuggerUrl")
                if ws_url:
                    return ws_url, launched_proc, launched_profile
    req = urllib.request.Request(
        "http://127.0.0.1:18800/json/new?" + urllib.parse.quote(page_url, safe=""),
        method="PUT",
    )
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        tab = json.loads(resp.read().decode("utf-8"))
    ws_url = tab.get("webSocketDebuggerUrl")
    if not ws_url:
        raise FetchError("Failed to create browser debugging page")
    return ws_url, launched_proc, launched_profile


def cdp_get_or_create_page_ws(page_url: str, timeout: int = 20) -> str:
    tabs = http_get_json("http://127.0.0.1:18800/json/list", timeout=timeout)
    if isinstance(tabs, list):
        for tab in tabs:
            tab_url = tab.get("url") or ""
            if tab_url == page_url or tab_url.startswith(page_url):
                ws_url = tab.get("webSocketDebuggerUrl")
                if ws_url:
                    return ws_url
    req = urllib.request.Request(
        "http://127.0.0.1:18800/json/new?" + urllib.parse.quote(page_url, safe=""),
        method="PUT",
    )
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        tab = json.loads(resp.read().decode("utf-8"))
    ws_url = tab.get("webSocketDebuggerUrl")
    if not ws_url:
        raise FetchError("无法创建浏览器调试页面")
    return ws_url


def cdp_call(ws: SimpleWebSocket, call_id: int, method: str, params: dict[str, Any] | None = None, wait_s: float = 10.0) -> dict[str, Any]:
    ws.send_json({"id": call_id, "method": method, "params": params or {}})
    deadline = time.time() + wait_s
    while time.time() < deadline:
        msg = ws.recv_json()
        if msg.get("id") == call_id:
            return msg
    raise FetchError(f"CDP 调用超时：{method}")


def extract_items_via_browser(page_url: str, count: int, wait_s: float = 8.0) -> tuple[dict[str, Any], list[dict[str, str]]]:
    use_headless = "ieeexplore.ieee.org" not in page_url.lower()
    ws_url, launched_proc, launched_profile = cdp_get_or_create_page_ws_auto(
        page_url,
        auto_launch=True,
        headless=use_headless,
    )
    ws = SimpleWebSocket(ws_url)
    ws.connect()
    try:
        cid = 1
        cdp_call(ws, cid, "Page.enable"); cid += 1
        cdp_call(ws, cid, "Runtime.enable"); cid += 1
        cdp_call(ws, cid, "Page.navigate", {"url": page_url}, wait_s=15.0); cid += 1
        time.sleep(12.0 if not use_headless else wait_s)
        consent_expr = """
(() => {
  const norm = (s) => (s || '').replace(/\s+/g, ' ').trim().toLowerCase();
  const buttons = Array.from(document.querySelectorAll('button, a, [role="button"], input[type="button"]'));
  const match = buttons.find((el) => {
    const text = norm(el.innerText || el.textContent || el.value || '');
    if (!text) return false;
    return (
      text === 'accept' ||
      text === 'i accept' ||
      text === 'agree' ||
      text === '接受' ||
      text.includes('accept all') ||
      text.includes('accept cookies') ||
      text.includes('接受全部')
    );
  });
  if (match) {
    match.click();
    return true;
  }
  return false;
})()
"""
        cdp_call(ws, cid, "Runtime.evaluate", {"expression": consent_expr, "returnByValue": True}, wait_s=10.0); cid += 1
        time.sleep(3.0)
        expr = f"""
(() => {{
  const norm = (s) => (s || '').replace(/\s+/g, ' ').trim();
  const banned = new Set(['editorial board', 'last update', 'copyright', 'about this issue', 'introduction']);
  const out = [];
  const seen = new Set();
  const linkSel = [
    '.article-list a',
    '.List-results-items a',
    '#results-blk .art-abs-url',
    'a[href*="/science/article/pii/"]',
    'a[href*="/document/"]',
    'a[href*="/abstract/document/"]',
    'a[href*="arnumber="]'
  ].join(',');
  const anchors = Array.from(document.querySelectorAll(linkSel));
  for (const link of anchors) {{
    const href = link.getAttribute('href') || '';
    const parentTag = (link.parentElement?.tagName || '').toLowerCase();
    const className = link.className || '';
    if (
      !href.includes('/science/article/pii/') &&
      !href.includes('/document/') &&
      !href.includes('arnumber=') &&
      !String(className).includes('art-abs-url')
    ) continue;
    if (parentTag && !['h2', 'h3'].includes(parentTag) && !String(className).includes('art-abs-url')) {{
      // On IEEE issue pages the title links live in h2/h3 or use art-abs-url.
      continue;
    }}
    const title = norm(link.innerText || link.textContent);
    if (!title || title.length < 8) continue;
    const key = title.toLowerCase();
    if (banned.has(key) || seen.has(key)) continue;
    seen.add(key);
    let authors = '';
    const container = link.closest('li, article, section, div') || document;
    const candidates = Array.from(container.querySelectorAll(
      'definition, dd, [class*="author"], [data-testid*="author"], .author, .authors, .doc-authors, .author-list'
    ));
    for (const el of candidates) {{
      const t = norm(el.innerText);
      if (!t) continue;
      if (/^Article\s+\d+/i.test(t)) continue;
      if (/^(View PDF|Article preview|Research article|Open access|Abstract only|Free access|PDF|HTML)$/i.test(t)) continue;
      if (t === title) continue;
      authors = t;
      break;
    }}
    out.push({{ title, authors }});
    if (out.length >= {max(1, count)}) break;
  }}
  const journal = norm(
    document.querySelector('meta[name="citation_journal_title"]')?.content ||
    document.querySelector('meta[property="og:site_name"]')?.content ||
    document.querySelector('h1')?.innerText
  ) || 'Journal';
  const volText = norm(Array.from(document.querySelectorAll('h2, div, span')).map(x => x.innerText).find(t => /(Volume|Vol\.)\s+\d+/i.test(norm(t))) || '');
  const m = volText.match(/(?:Volume|Vol\.)\s+(\d+)/i);
  return {{
    journal,
    volume: m ? m[1] : 'Unknown',
    items: out,
    pageTitle: document.title || '',
    bodySample: norm(document.body?.innerText || '').slice(0, 800)
  }};
}})()
"""
        for _ in range(12):
            scroll_expr = """
(() => {
  window.scrollTo(0, document.body.scrollHeight);
  const more = Array.from(document.querySelectorAll('button, a, [role="button"]')).find((el) => {
    const text = ((el.innerText || el.textContent || '')).replace(/\s+/g, ' ').trim().toLowerCase();
    return text.includes('show more') || text.includes('load more') || text.includes('more');
  });
  if (more) more.click();
  return {
    height: document.body.scrollHeight,
    itemsHint: document.querySelectorAll('h2 a, h3 a, .art-abs-url, a[href*="/document/"], a[href*="arnumber="]').length
  };
})()
"""
            cdp_call(ws, cid, "Runtime.evaluate", {"expression": scroll_expr, "returnByValue": True}, wait_s=10.0); cid += 1
            time.sleep(1.5)
            result = cdp_call(ws, cid, "Runtime.evaluate", {"expression": expr, "returnByValue": True}, wait_s=20.0); cid += 1
            value = (((result.get("result") or {}).get("result") or {}).get("value"))
            if value and len(value.get("items") or []) >= count:
                break
        else:
            result = cdp_call(ws, cid, "Runtime.evaluate", {"expression": expr, "returnByValue": True}, wait_s=20.0)
        value = (((result.get("result") or {}).get("result") or {}).get("value"))
        if not value or not value.get("items"):
            page_title = (value or {}).get("pageTitle", "")
            body_sample = (value or {}).get("bodySample", "")
            raise FetchError(
                f"Browser extraction failed for URL: {page_url} | title={page_title!r} | sample={body_sample!r}"
            )
        if not value or not value.get("items"):
            raise FetchError(f"浏览器模式未能提取网页条目：{page_url}")
        meta = {"title": value.get("journal") or "Journal", "volume": value.get("volume") or "Unknown", "source_url": page_url}
        items = [
            {"title": normalize_title(it.get("title", "")), "authors": normalize_title(it.get("authors", "")) or "Unknown", "title_zh": "", "source": "BrowserDOM"}
            for it in value.get("items", [])
            if normalize_title(it.get("title", ""))
        ]
        return meta, items[:count]
    finally:
        ws.close()
        if launched_proc is not None:
            try:
                launched_proc.terminate()
                launched_proc.wait(timeout=5)
            except Exception:
                pass
        if launched_profile:
            try:
                import shutil
                shutil.rmtree(launched_profile, ignore_errors=True)
            except Exception:
                pass


def extract_items_from_url_page(page_url: str, count: int, strict_order: bool = False) -> tuple[dict[str, Any], list[dict[str, str]]]:
    if strict_order:
        return extract_items_via_browser(page_url, count)

    if "ieeexplore.ieee.org" in page_url.lower() and "isnumber=" in page_url.lower() and "punumber=" in page_url.lower():
        try:
            return extract_ieee_items_via_api(page_url, count)
        except Exception:
            pass

    html_text = ""
    try:
        html_text = http_get_text(page_url)
    except Exception:
        html_text = ""

    journal_title = ""
    volume = ""
    if html_text:
        journal_title, volume = guess_journal_and_volume_from_html(page_url, html_text)
    inferred_journal, inferred_volume = infer_journal_volume_from_url(page_url)
    if (
        not journal_title
        or journal_title == "Journal"
        or journal_title.lower() == "sciencedirect"
        or journal_title.startswith("http://")
        or journal_title.startswith("https://")
    ):
        journal_title = inferred_journal
    if not volume or volume == "Unknown":
        volume = inferred_volume

    candidates: list[dict[str, str]] = []

    if html_text:
        # 1) JSON-LD blocks
        for m in re.finditer(r'<script[^>]+type=["\']application/ld\+json["\'][^>]*>(.*?)</script>', html_text, flags=re.I | re.S):
            raw = m.group(1).strip()
            if not raw:
                continue
            try:
                payload = json.loads(raw)
            except Exception:
                continue
            _walk_json_for_article_like_entries(payload, candidates)

        # 2) IEEE Xplore issue/search pages
        if len(candidates) < count and "ieeexplore.ieee.org" in page_url.lower():
            candidates.extend(extract_ieee_issue_entries_from_html(html_text, count - len(candidates)))

        # 3) ScienceDirect-ish title spans as fallback
        if len(candidates) < count:
            titles = [strip_tags(x) for x in SCIENCEDIRECT_TITLE_BLOCK.findall(html_text)]
            titles = [t for t in titles if t and len(t) > 8]
            for title in titles:
                candidates.append({"title": title, "authors": "Unknown", "title_zh": "", "source": "URL"})

        # 4) Generic citation meta tags for single-article pages
        if len(candidates) < count:
            meta_title = extract_meta_content(html_text, "citation_title")
            if meta_title:
                meta_authors = re.findall(r'<meta[^>]+name=["\']citation_author["\'][^>]+content=["\']([^"\']+)["\']', html_text, flags=re.I)
                candidates.append({
                    "title": normalize_title(meta_title),
                    "authors": ", ".join(html.unescape(x).strip() for x in meta_authors) or "Unknown",
                    "title_zh": "",
                    "source": "URL",
                })

    # Deduplicate and trim obvious non-paper entries
    results: list[dict[str, str]] = []
    seen: set[str] = set()
    banned = {"editorial board", "last update", "copyright", "about this issue", "introduction"}
    for item in candidates:
        title = normalize_title(item.get("title", ""))
        if not title:
            continue
        if title.lower() in banned:
            continue
        if title in seen:
            continue
        seen.add(title)
        results.append({
            "title": title,
            "authors": item.get("authors") or "Unknown",
            "title_zh": item.get("title_zh") or "",
            "source": item.get("source") or "URL",
        })
        if len(results) >= count:
            break

    browser_error: Exception | None = None
    if not results:
        try:
            return extract_items_via_browser(page_url, count)
        except Exception as exc:
            browser_error = exc

    if not results and journal_title and volume and volume != "Unknown":
        crossref_meta, crossref_items = fetch_volume_works(journal_title, volume, count)
        crossref_meta["source_url"] = page_url
        crossref_meta["title"] = crossref_meta.get("title") or journal_title
        return crossref_meta, crossref_items

    if not results and browser_error is not None:
        raise FetchError(f"Unable to extract paper entries from URL: {page_url} | browser_fallback={browser_error}")

    if not results:
        raise FetchError(f"未能从链接提取论文条目：{page_url}")

    if not results and browser_error is not None:
        raise FetchError(f"Unable to extract paper entries from URL: {page_url} | browser_fallback={browser_error}")

    meta = {
        "title": journal_title,
        "source_url": page_url,
        "volume": volume,
    }
    return meta, results


def fetch_volume_works(journal_title: str, volume: str, count: int, per_page: int = 100) -> tuple[dict[str, Any], list[dict[str, str]]]:
    journal = best_journal_match(journal_title)
    issn_list = journal.get("ISSN", [])
    if not issn_list:
        raise FetchError(f"期刊 {journal.get('title') or journal_title} 没有可用 ISSN")
    issn = issn_list[0]

    rows_needed = max(count * 8, per_page)
    rows_needed = min(rows_needed, 1000)
    url = f"{CROSSREF_API}/journals/{issn}/works?rows={rows_needed}&sort=published&order=desc"
    data = http_get_json(url)
    items = data.get("message", {}).get("items", [])

    seen = set()
    results: list[dict[str, str]] = []
    target_volume = str(volume).strip()
    for item in items:
        item_volume = str(item.get("volume", "")).strip()
        if item_volume != target_volume:
            continue
        titles = item.get("title") or []
        if not titles:
            continue
        title = normalize_title(titles[0])
        if not title or title in seen:
            continue
        seen.add(title)
        authors = parse_authors(item.get("author") or [])
        results.append({
            "title": title,
            "authors": authors,
            "title_zh": "",
            "doi": item.get("DOI", ""),
            "issued": json.dumps(item.get("issued", {}), ensure_ascii=False),
            "source": "Crossref",
        })
        if len(results) >= count:
            break

    if not results:
        raise FetchError(f"期刊 {journal.get('title') or journal_title} 的第 {volume} 卷没有查到论文")

    journal = dict(journal)
    journal["volume"] = str(volume)
    return journal, results


def translate_text_google(text: str, target_lang: str = "zh-CN", source_lang: str = "auto") -> str:
    q = urllib.parse.quote(text)
    url = (
        "https://translate.googleapis.com/translate_a/single"
        f"?client=gtx&sl={source_lang}&tl={target_lang}&dt=t&q={q}"
    )
    raw = http_get_text(url)
    data = json.loads(raw)
    pieces = []
    for part in data[0]:
        if part and part[0]:
            pieces.append(part[0])
    return "".join(pieces).strip()


def enrich_translations(items: list[dict[str, str]], sleep_s: float = 0.2) -> None:
    for i, item in enumerate(items, 1):
        title = item["title"]
        try:
            item["title_zh"] = translate_text_google(title)
        except Exception:
            item["title_zh"] = title
        if i != len(items):
            time.sleep(sleep_s)


def add_bg(prs: Presentation, slide, dark: bool = False):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY if dark else LIGHT
    if not dark:
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.35))
        band.fill.solid()
        band.fill.fore_color.rgb = NAVY
        band.line.fill.background()
        footer = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.28), prs.slide_width, Inches(0.28)
        )
        footer.fill.solid()
        footer.fill.fore_color.rgb = RGBColor(231, 236, 242)
        footer.line.fill.background()


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK, name="Aptos", align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = name
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    return tb


def estimate_lines(text: str, chars_per_line: int) -> int:
    return max(1, math.ceil(len(text or "") / chars_per_line))


def fits_item(item: dict[str, str], layout: dict[str, float]) -> bool:
    title_lines = estimate_lines(item["title"], layout["title_chars"])
    author_lines = estimate_lines("作者：" + item["authors"], layout["author_chars"])
    zh_lines = estimate_lines("中文：" + item["title_zh"], layout["zh_chars"])
    return title_lines <= 2 and author_lines <= 1 and zh_lines <= 1


def fits_page(items: list[dict[str, str]], layout: dict[str, float]) -> bool:
    if len(items) * layout["card_h"] + (len(items) - 1) * CARD_GAP > CONTENT_HEIGHT:
        return False
    return all(fits_item(item, layout) for item in items)


def paginate(items: list[dict[str, str]]):
    pages = []
    i = 0
    n = len(items)
    while i < n:
        chunk5 = items[i : i + 5]
        if len(chunk5) == 5 and fits_page(chunk5, LAYOUT5):
            pages.append((chunk5, 5))
            i += 5
            continue
        chunk4 = items[i : i + 4]
        pages.append((chunk4, 4))
        i += 4
    return pages


def add_cover(prs: Presentation, journal_title: str, volume: str, count: int, source_name: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(prs, slide, dark=True)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(0.9), Inches(0.18), Inches(5.6))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(120, 167, 194)
    accent.line.fill.background()
    add_text(slide, 1.1, 1.0, 10.9, 1.2, f"{journal_title} 第{volume}卷前{count}篇论文信息汇编", 28, True, WHITE)
    add_text(slide, 1.1, 2.0, 10.7, 0.8, "Titles · Authors · Chinese Translations", 18, False, RGBColor(210, 221, 233))
    add_text(
        slide,
        1.1,
        3.0,
        11.0,
        2.5,
        f"来源：{source_name}\n内容范围：按返回顺序整理前 {count} 篇论文\n编排方式：统一字号；优先每页5篇，放不下时自动改为4篇",
        20,
        False,
        WHITE,
    )
    add_text(slide, 1.1, 6.45, 11, 0.3, "Generated locally with python-pptx", 11, False, RGBColor(195, 205, 218))


def add_overview(prs: Presentation, total_pages: int, count5: int, count4: int, journal_title: str, volume: str, item_count: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(prs, slide, dark=False)
    add_text(slide, 0.7, 0.55, 8, 0.5, "内容说明", 24, True, NAVY)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.25), Inches(11.8), Inches(4.9))
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = ACCENT
    items = [
        f"本汇编收录 {journal_title} 第 {volume} 卷前 {item_count} 篇论文。",
        "每篇给出：英文标题、作者、中文标题翻译。",
        "正文统一字号，不对单条标题做单独缩放。",
        "分页规则：先尝试每页 5 篇；若 5 篇排版会拥挤，则自动改为 4 篇。",
        f"当前正文共 {total_pages} 页，其中 5 篇页 {count5} 页，4 篇页 {count4} 页。",
    ]
    y = 1.55
    for i, it in enumerate(items, 1):
        add_text(slide, 1.2, y, 10.8, 0.55, f"{i}. {it}", 19, False, DARK)
        y += 0.8
    add_text(slide, 0.8, 6.8, 12, 0.25, "注：元数据来自 Crossref；中文标题为脚本自动翻译，仅供参考。", 10.5, False, MID)


def add_entry_slide(prs: Presentation, items: list[dict[str, str]], mode: int, start_idx: int, end_idx: int, journal_title: str, volume: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(prs, slide, dark=False)
    add_text(slide, 0.7, 0.52, 6, 0.45, f"论文条目 {start_idx}–{end_idx}", 24, True, NAVY)
    add_text(slide, 9.1, 0.56, 3.5, 0.3, f"{journal_title} Vol.{volume}", 11, False, TEAL, align=PP_ALIGN.RIGHT)

    layout = LAYOUT5 if mode == 5 else LAYOUT4
    y = CONTENT_TOP
    for idx, item in enumerate(items, start_idx):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(CARD_X), Inches(y), Inches(CARD_W), Inches(layout["card_h"]))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = ACCENT

        num = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(NUM_X), Inches(y + 0.12), Inches(0.55), Inches(0.42))
        num.fill.solid()
        num.fill.fore_color.rgb = TEAL
        num.line.fill.background()
        add_text(slide, NUM_X, y + 0.14, 0.55, 0.25, str(idx), 16, True, WHITE, align=PP_ALIGN.CENTER)

        add_text(slide, LEFT, y + layout["title_y"], TEXT_W, layout["title_h"], item["title"], TITLE_FONT, True, DARK)
        add_text(slide, LEFT, y + layout["author_y"], TEXT_W, layout["author_h"], "作者：" + item["authors"], AUTHOR_FONT, False, MID)
        add_text(slide, LEFT, y + layout["zh_y"], TEXT_W, layout["zh_h"], "中文：" + (item["title_zh"] or item["title"]), ZH_FONT, False, TEAL)
        y += layout["card_h"] + CARD_GAP


def build_ppt(items: list[dict[str, str]], journal_title: str, volume: str, out_path: Path, source_name: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_cover(prs, journal_title, volume, len(items), source_name)
    pages = paginate(items)
    count5 = sum(1 for _, mode in pages if mode == 5)
    count4 = sum(1 for _, mode in pages if mode == 4)
    add_overview(prs, len(pages), count5, count4, journal_title, volume, len(items))

    start = 1
    for chunk, mode in pages:
        end = start + len(chunk) - 1
        add_entry_slide(prs, chunk, mode, start, end, journal_title, volume)
        start = end + 1

    prs.save(str(out_path))
    return len(pages) + 2, len(pages), count5, count4


def default_output_name(journal_title: str, volume: str, count: int) -> str:
    base = slugify(journal_title)
    if volume and volume != "Unknown":
        return f"{base}_vol{volume}_first{count}_title_author_cn.pptx"
    return f"{base}_first{count}_title_author_cn.pptx"


def default_json_name(journal_title: str, volume: str, count: int) -> str:
    base = slugify(journal_title)
    if volume and volume != "Unknown":
        return f"{base}_vol{volume}_first{count}.json"
    return f"{base}_first{count}.json"


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Generate a PPT for the first N papers in a journal volume or URL page")
    p.add_argument("--journal", help="Journal title, e.g. Automatica")
    p.add_argument("--volume", help="Volume number, e.g. 185")
    p.add_argument("--url", help="Issue/article-list page URL. If provided, the script extracts titles from the page directly")
    p.add_argument("--count", required=True, type=int, help="Number of paper titles to include")
    p.add_argument("--strict-order", action="store_true", help="For URL mode, extract titles in the exact rendered page order via local Chrome DevTools")
    p.add_argument("--output", help="Output PPTX path")
    p.add_argument("--json", dest="json_path", help="Optional path to save fetched metadata JSON")
    p.add_argument("--no-translate", action="store_true", help="Do not auto-translate titles to Chinese")
    p.add_argument("--translate-sleep", type=float, default=0.2, help="Delay between translations in seconds")
    args = p.parse_args()
    if not args.url and not (args.journal and args.volume):
        p.error("必须提供 --url，或者同时提供 --journal 和 --volume")
    return args


def main() -> int:
    args = parse_args()
    if args.count <= 0:
        print("count 必须大于 0", file=sys.stderr)
        return 2

    base_journal = args.journal or "journal"
    base_volume = args.volume or "unknown"

    if args.url:
        source_meta, items = extract_items_from_url_page(args.url, args.count, strict_order=args.strict_order)
        journal_title = source_meta.get("title") or base_journal
        volume = str(source_meta.get("volume") or base_volume)
        source_name = source_meta.get("source_url") or args.url
        payload_source = "BrowserDOM" if args.strict_order else "URL"
        payload_extra = {"source_url": args.url, "strict_order": args.strict_order}
    else:
        source_meta, items = fetch_volume_works(args.journal, args.volume, args.count)
        journal_title = source_meta.get("title") or args.journal
        volume = str(args.volume)
        source_name = f"Crossref / {journal_title}"
        payload_source = "Crossref"
        payload_extra = {"issn": source_meta.get("ISSN", [])}

    out_path = Path(args.output) if args.output else Path(default_output_name(journal_title, volume, args.count))
    json_path = Path(args.json_path) if args.json_path else Path(default_json_name(journal_title, volume, args.count))

    if not args.no_translate:
        enrich_translations(items, sleep_s=args.translate_sleep)
    else:
        for item in items:
            item["title_zh"] = item["title"]

    payload = {
        "journal_query": args.journal,
        "journal_title": journal_title,
        "volume": volume,
        "count": len(items),
        "source": payload_source,
        **payload_extra,
        "items": items,
    }
    json_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")

    slides_total, content_slides, count5, count4 = build_ppt(
        items,
        journal_title,
        volume,
        out_path,
        source_name,
    )

    print(f"JSON saved: {json_path}")
    print(
        f"PPT created: {out_path} | slides={slides_total} | content_slides={content_slides} | "
        f"5-item pages={count5} | 4-item pages={count4}"
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
